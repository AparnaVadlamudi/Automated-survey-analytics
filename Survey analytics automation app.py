import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from wordcloud import WordCloud
from io import BytesIO
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches, Pt
import numpy as np
from sklearn.feature_extraction.text import CountVectorizer
from matplotlib import cm

# Set the maximum sample size for large datasets
MAX_SAMPLE_SIZE = 100

st.title('Survey Analytics Automation')

# Add an introduction
st.markdown("""
Welcome to the **Survey Analytics Automation** app!
This tool allows you to upload your survey data (CSV or Excel format) and automatically generate insightful visualizations and analyses.
You can customize data types, select chart types, and export the results in PDF or PPTX format.
""")

uploaded_file = st.file_uploader("Choose a CSV or Excel file", type=["csv", "xlsx"])

# Global Color Theme Selection
st.subheader("Select Color Theme for Visualizations")
color_theme_options = ['Default', 'Magma', 'Inferno', 'Viridis', 'Cividis', 'Pastel', 'Bright', 'Dark']
selected_theme_name = st.selectbox("Select Color Theme", color_theme_options, index=0, key='color_theme_select')

# Custom Default Colors
custom_default_colors = ['#5E4FA2', '#3288BD', '#66C2A5', '#ABDDA4', '#E6F598', 
                         '#FEE08B', '#FDAE61', '#F46D43', '#D53E4F', '#9E0142']

# Mapping of selected theme to colormaps
colormap_mapping = {
    'Default': custom_default_colors,
    'Magma': 'magma',
    'Inferno': 'inferno',
    'Viridis': 'viridis',
    'Cividis': 'cividis',
    'Pastel': 'Pastel1',
    'Bright': 'Set1',
    'Dark': 'Dark2'
}

# Set selected color palette
if selected_theme_name == 'Default':
    selected_colormap = 'crest'  # For word cloud when 'Default' is selected
    selected_palette = custom_default_colors
else:
    selected_colormap = colormap_mapping[selected_theme_name]
    selected_palette = sns.color_palette(selected_colormap).as_hex()

# Set seaborn palette for plots
sns.set_palette(selected_palette)

def get_colors_from_cmap(cmap_name, num_colors):
    cmap = cm.get_cmap(cmap_name)
    colors = [cmap(i) for i in np.linspace(0, 1, num_colors)]
    return colors

import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from io import BytesIO
import textwrap

def export_to_pdf(charts, insights, raw_texts):
    pdf_io = BytesIO()
    with PdfPages(pdf_io) as pdf:
        for col in charts:
            if charts[col]:
                fig = charts[col]
                
                # Determine if the chart is a word cloud based on the title or content
                is_word_cloud = 'Word Cloud' in col or 'wordcloud' in col.lower()
                
                # Set figure size to A4 portrait (8.27 x 11.69 inches)
                fig.set_size_inches(8.27, 6)
                
                # Adjust subplot to allocate space for x-axis title and insights
                if is_word_cloud:
                    fig.subplots_adjust(top=0.80, bottom=0.25)  # More space for word clouds
                else:
                    fig.subplots_adjust(top=0.85, bottom=0.30)  # Standard adjustment
                
                # Wrap the main title if it's too long
                wrapped_title = "\n".join(textwrap.wrap(col, width=60))
                fig.suptitle(wrapped_title, fontsize=16, y=0.95, ha='center')
                
                # Adjust x-axis title to take full width and align it properly
                for ax in fig.axes:
                    if ax.get_xlabel():
                        xlabel = ax.get_xlabel()
                        wrapped_xlabel = "\n".join(textwrap.wrap(xlabel, width=80))  # Wider wrap for full-width alignment
                        ax.set_xlabel(wrapped_xlabel, fontsize=12, labelpad=10, ha='center')
                    
                    # Rotate x-axis labels for better visibility
                    ax.tick_params(axis='x', rotation=45, labelsize=10)
                
                # Position insights closer to the word cloud image
                insight_text = insights.get(col, "")
                if insight_text:
                    wrapped_insight = "\n".join(textwrap.wrap(insight_text, width=100))
                    if is_word_cloud:
                        # Increase font size for word cloud insights and position them closer to the image
                        fig.text(0.5, 0.15, f"Insights: {wrapped_insight}", fontsize=14, ha='center', va='top', wrap=True)
                    else:
                        # Standard font size for other charts
                        fig.text(0.5, 0.02, f"Insights: {wrapped_insight}", fontsize=12, ha='center', va='bottom', wrap=True)
                
                pdf.savefig(fig)
                plt.close(fig)

            elif raw_texts[col]:
                # Handle raw responses
                fig, ax = plt.subplots(figsize=(8.27, 11.69))  # A4 portrait size
                ax.axis('off')
                ax.axis('tight')
                
                # Add title with padding
                wrapped_title = "\n".join(textwrap.wrap(f'Raw Responses for {col}', width=60))
                fig.suptitle(wrapped_title, fontsize=16, y=0.90, ha='center')
                
                # Prepare data for the table
                data = [[response] for response in raw_texts[col].split('\n\n')]
                
                # Create table and set word wrapping
                table = ax.table(cellText=data, colLabels=[col], loc='upper center', cellLoc='left')
                table.auto_set_font_size(False)
                table.set_fontsize(10)
                table.scale(1, 1.5)
                
                # Add padding to table cells and enable word wrapping
                for key, cell in table.get_celld().items():
                    cell.set_text_props(wrap=True)
                    cell.set_linewidth(0.5)  # Ensure lines are visible
                    cell.pad = 0.5  # Add padding for better spacing

                # Adjust layout to prevent overlapping
                plt.tight_layout(rect=[0, 0.05, 1, 0.90])  # [left, bottom, right, top]
                
                pdf.savefig(fig)
                plt.close(fig)
    
    pdf_io.seek(0)
    return pdf_io.getvalue()


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from io import BytesIO

def export_to_pptx(charts, insights, raw_texts):
    ppt_io = BytesIO()
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for col in charts:
        # Use 'Blank' layout for full control
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Adjusted title text box size and position
        title_left = Inches(0.5)
        title_top = Inches(0.2)
        title_width = Inches(9)
        title_height = Inches(2.0)  # Increased height to accommodate multiple lines
        
        title_placeholder = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_tf = title_placeholder.text_frame
        title_tf.word_wrap = True  # Enable word wrap for the title
        title_tf.auto_size = MSO_AUTO_SIZE.NONE  # Prevent auto-sizing
        
        title_p = title_tf.paragraphs[0]
        title_p.text = col
        title_p.font.size = Pt(26)
        title_p.font.bold = True
        
        # Estimate the number of lines in the title
        # This is a rough estimation based on character count and title width
        max_chars_per_line = 60  # Adjust as necessary
        num_lines = (len(col) // max_chars_per_line) + 1
        
        # Calculate the actual height occupied by the title
        # Assuming each line is approximately 0.3 inches in height
        actual_title_height = Inches(0.4 * num_lines)
        
        # Update the top position for the content based on title height
        content_top = title_top + actual_title_height + Inches(0.9)  # Add a small margin
        
        if charts[col]:
            fig = charts[col]
            fig.set_size_inches(7.5, 4.0)  # Adjusted size to fit within slide
            fig.subplots_adjust(top=0.85, bottom=0.2)
    
            img_io = BytesIO()
            fig.savefig(img_io, format='png', bbox_inches='tight')
            img_io.seek(0)
    
            # Adjusted positions to prevent overlap
            left = Inches(0.5)
            image_height = Inches(4.5)  # Fixed image height
            slide.shapes.add_picture(img_io, left, content_top, width=Inches(9), height=image_height)
            plt.close(fig)
    
            # Add insights below the chart
            insight_text = insights.get(col, "")
            if insight_text:
                insight_top = content_top + image_height + Inches(0.1)  # Position below the image
                txBox = slide.shapes.add_textbox(Inches(0.5), insight_top, Inches(9), Inches(1.5))
                tf = txBox.text_frame
                tf.word_wrap = True  # Enable word wrapping
                p = tf.paragraphs[0]
                p.text = f"Insights: {insight_text}"
                p.font.size = Pt(12)
                p.font.bold = False
        elif raw_texts[col]:
            # Handle raw responses
            raw_responses = raw_texts[col].split('\n\n')
            rows = len(raw_responses) + 1
            cols_table = 1
            left = Inches(0.25)
            top = content_top  # Start below the title
            width = Inches(10)
            height = prs.slide_height - top - Inches(0.5)  # Adjust height to fit within slide
    
            table_shape = slide.shapes.add_table(rows, cols_table, left, top, width, height)
            table = table_shape.table
            table.columns[0].width = Inches(9)
            cell = table.cell(0, 0)
            cell.text = col
    
            # Set consistent font size for header
            header_run = cell.text_frame.paragraphs[0].runs[0]
            header_run.font.size = Pt(14)
            header_run.font.bold = True
    
            for i, response in enumerate(raw_responses, start=1):
                cell = table.cell(i, 0)
                cell.text = response
                # Set consistent font size and enable word wrap
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(12)
                        run.font.name = 'Arial'
                cell.text_frame.word_wrap = True
    
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io.getvalue()



if uploaded_file is not None:
    try:
        # Determine the file type
        if uploaded_file.name.endswith('.csv'):
            df = pd.read_csv(uploaded_file)
        elif uploaded_file.name.endswith('.xlsx'):
            df = pd.read_excel(uploaded_file)
        else:
            st.error("Unsupported file format. Please upload a CSV or Excel file.")
            st.stop()

        # Prompt user to select columns to exclude
        st.subheader("Select Columns to Exclude from Analysis")
        st.write("You may exclude columns such as IDs, names, timestamps, etc., that are not useful for analysis.")
        exclude_cols = st.multiselect("Select columns to exclude:", df.columns, key='exclude_cols')

        # Filter out the columns to be excluded
        df_filtered = df.drop(columns=exclude_cols)

        # Sample the data if necessary
        if len(df_filtered) > MAX_SAMPLE_SIZE:
            st.info(f"Dataset is large ({len(df_filtered)} rows). Sampling {MAX_SAMPLE_SIZE} rows for initial analysis.")
            df_sample = df_filtered.sample(n=MAX_SAMPLE_SIZE, random_state=42)
        else:
            df_sample = df_filtered

        # Data type detection with improved Likert/Rating detection
        column_data_types = {}
        for col in df_sample.columns:
            non_null_values = df_sample[col].dropna()
            if non_null_values.empty:
                column_data_types[col] = 'Other'
                continue

            unique_values = non_null_values.unique()
            num_unique = len(unique_values)
            non_null_count = len(non_null_values)

            if pd.api.types.is_numeric_dtype(non_null_values):
                if (non_null_values.astype(float) == non_null_values.astype(int)).all():
                    min_value = unique_values.min()
                    max_value = unique_values.max()
                    if (min_value >= 1 and max_value <= 10) and (num_unique <= 10):
                        column_data_types[col] = 'Likert/Rating'
                    else:
                        column_data_types[col] = 'Numeric'
                else:
                    column_data_types[col] = 'Numeric'
            elif pd.api.types.is_string_dtype(non_null_values):
                avg_length = non_null_values.str.len().mean()
                if avg_length <= 20:
                    if num_unique <= 5:
                        column_data_types[col] = 'Likert/Rating'
                    elif num_unique / non_null_count < 0.1:
                        column_data_types[col] = 'Categorical'
                    else:
                        column_data_types[col] = 'Text'
                else:
                    column_data_types[col] = 'Text'
            else:
                column_data_types[col] = 'Other'

        # Allow user to override data types
        st.subheader("Data Type Confirmation and Override")
        st.write("The default selection has been made based on the most suitable data type derived from the data entries. You can modify these selections if needed.")

        # Use columns to arrange data type selectors side by side
        num_columns = 2  # Adjusted to 2 columns
        cols = st.columns(num_columns)
        for idx, col in enumerate(df_sample.columns):
            detected_type = column_data_types[col]
            with cols[idx % num_columns]:
                data_type = st.selectbox(f"Data type for '{col}'", ['Numeric', 'Categorical', 'Text', 'Likert/Rating', 'Other'],
                                         index=['Numeric', 'Categorical', 'Text', 'Likert/Rating', 'Other'].index(detected_type),
                                         key=f"data_type_{col}")
                column_data_types[col] = data_type

        charts = {}
        insights = {}
        raw_texts = {}

        for col in df_sample.columns:
            with st.expander(f"{col}", expanded=True):
                # Check if the column has any data to analyze
                if df[col].dropna().empty:
                    st.warning(f"No data available for column '{col}'. Please select another column.")
                    continue

                data_type = column_data_types[col]
                st.write(f"Using data type: {data_type}")

                unique_values = df[col].nunique()

                # Generate colors from the selected colormap
                colors = get_colors_from_cmap(selected_colormap, max(unique_values, 3))

                # Simplify chart type selection
                chart_cols = st.columns([1, 3])  # Adjust column widths
                if data_type == 'Numeric':
                    default_chart = 'Histogram'
                    chart_options = ['Histogram', 'Boxplot', 'Violin Plot', 'Scatter Plot']
                elif data_type == 'Categorical':
                    default_chart = 'Bar Chart'
                    chart_options = ['Bar Chart', 'Pie Chart', 'Count Plot']
                elif data_type == 'Likert/Rating':
                    default_chart = 'Bar Chart'
                    chart_options = ['Bar Chart', 'Pie Chart', 'Stacked Bar Chart']
                elif data_type == 'Text':
                    default_chart = 'Word Cloud'
                    chart_options = ['Word Cloud', 'Raw Responses']
                else:
                    default_chart = 'Bar Chart'
                    chart_options = ['Bar Chart', 'Pie Chart']

                with chart_cols[0]:
                    chart_type = st.selectbox("Select Chart Type", chart_options,
                                              index=chart_options.index(default_chart),
                                              key=f"chart_type_{col}")
                with chart_cols[1]:
                    pass  # Empty column to adjust layout

                fig = None
                raw_text = None  # To handle raw responses

                if data_type == 'Numeric':
                    if chart_type == 'Histogram':
                        fig, ax = plt.subplots(figsize=(8, 5))
                        sns.histplot(df[col].dropna(), kde=True, ax=ax, color=colors[0])
                        ax.set_xlabel(col, fontsize=12)
                        ax.set_ylabel('Frequency', fontsize=12)
                        st.pyplot(fig)
                    elif chart_type == 'Boxplot':
                        fig, ax = plt.subplots(figsize=(8, 5))
                        sns.boxplot(x=df[col], ax=ax, color=colors[0])
                        ax.set_xlabel(col, fontsize=12)
                        st.pyplot(fig)
                    elif chart_type == 'Violin Plot':
                        fig, ax = plt.subplots(figsize=(8, 5))
                        sns.violinplot(x=df[col], ax=ax, color=colors[0])
                        ax.set_xlabel(col, fontsize=12)
                        st.pyplot(fig)
                    elif chart_type == 'Scatter Plot':
                        numeric_columns = [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c]) and c != col]
                        if numeric_columns:
                            selected_col = st.selectbox("Select another numeric column for scatter plot", numeric_columns, key=f"{col}_scatter")
                            fig, ax = plt.subplots(figsize=(8, 5))
                            ax.scatter(df[col], df[selected_col], color=colors[0])
                            ax.set_xlabel(col, fontsize=12)
                            ax.set_ylabel(selected_col, fontsize=12)
                            st.pyplot(fig)
                        else:
                            st.write("No other numeric columns available for scatter plot.")
                    insights[col] = ""
                    raw_texts[col] = ""

                elif data_type == 'Categorical':
                    if chart_type == 'Bar Chart':
                        fig, ax = plt.subplots(figsize=(8, 5))
                        df[col].value_counts().plot(kind='bar', ax=ax, color=colors)
                        ax.set_xlabel(col, fontsize=12)
                        ax.set_ylabel('Count of Responses', fontsize=12)
                        ax.tick_params(axis='x', rotation=0)
                        st.pyplot(fig)
                    elif chart_type == 'Pie Chart':
                        fig, ax = plt.subplots(figsize=(8, 8))  # Adjusted size
                        df[col].value_counts().plot(kind='pie', ax=ax, autopct='%1.1f%%', colors=colors, textprops={'fontsize': 10})
                        ax.set_ylabel('')
                        plt.legend(title=col, bbox_to_anchor=(1, 0.5), loc='center left', fontsize=10)
                        plt.tight_layout()
                        st.pyplot(fig)
                    elif chart_type == 'Count Plot':
                        fig, ax = plt.subplots(figsize=(8, 5))
                        sns.countplot(y=col, data=df, ax=ax, palette=colors)
                        ax.set_xlabel('Count of Responses', fontsize=12)
                        ax.set_ylabel(col, fontsize=12)
                        st.pyplot(fig)
                    insights[col] = ""
                    raw_texts[col] = ""

                elif data_type == 'Likert/Rating':
                    if chart_type == 'Bar Chart':
                        fig, ax = plt.subplots(figsize=(8, 5))
                        value_counts = df[col].value_counts(sort=False).sort_index()
                        value_counts.plot(kind='bar', ax=ax, color=colors)
                        ax.set_xlabel(col, fontsize=12)
                        ax.set_ylabel('Count of Responses', fontsize=12)
                        ax.tick_params(axis='x', rotation=0)

                        total = value_counts.sum()
                        for p in ax.patches:
                            percentage = f'{(p.get_height() / total) * 100:.1f}%'
                            ax.annotate(
                                percentage,
                                (p.get_x() + p.get_width() / 2., p.get_height() / 2),
                                ha='center', va='center',
                                color='white', fontsize=10, fontweight='bold'
                            )

                        plt.tight_layout()
                        st.pyplot(fig)

                        if df[col].dtype == 'object':
                            df[col] = pd.to_numeric(df[col], errors='coerce')
                        if df[col].notnull().any():
                            average_rating = df[col].mean()
                            top_ratings = df[col][df[col] >= 4].count()
                            top_ratings_percentage = (top_ratings / total) * 100
                            insights[col] = f"The average rating was {average_rating:.2f}. {top_ratings_percentage:.1f}% of respondents rated '{col}' as 4 or higher."
                            st.markdown(f"**Insight**: {insights[col]}")
                        else:
                            insights[col] = ""
                            st.markdown(f"**Insight**: No valid numeric ratings available for this question.")
                        raw_texts[col] = ""

                    elif chart_type == 'Pie Chart':
                        fig, ax = plt.subplots(figsize=(8, 8))  # Adjusted size
                        value_counts = df[col].value_counts(sort=False).sort_index()
                        value_counts.plot(kind='pie', ax=ax, autopct='%1.1f%%', colors=colors, textprops={'fontsize': 10})
                        ax.set_ylabel('')
                        plt.legend(title=col, bbox_to_anchor=(1, 0.5), loc='center left', fontsize=10)
                        plt.tight_layout()
                        st.pyplot(fig)

                        if df[col].dtype == 'object':
                            df[col] = pd.to_numeric(df[col], errors='coerce')
                        if df[col].notnull().any():
                            average_rating = df[col].mean()
                            top_ratings = df[col][df[col] >= 4].count()
                            total = df[col].count()
                            top_ratings_percentage = (top_ratings / total) * 100
                            insights[col] = f"The average rating was {average_rating:.2f}. {top_ratings_percentage:.1f}% of respondents rated '{col}' as 4 or higher."
                            st.markdown(f"**Insight**: {insights[col]}")
                        else:
                            insights[col] = ""
                            st.markdown(f"**Insight**: No valid numeric ratings available for this question.")
                        raw_texts[col] = ""

                    elif chart_type == 'Stacked Bar Chart':
                        categorical_columns = [c for c in df.columns if c != col and df[c].dtype == 'object']
                        if categorical_columns:
                            stack_by = st.selectbox("Select a categorical column to stack by", categorical_columns, key=f"{col}_stack")
                            stacked_data = df.groupby([stack_by, col]).size().unstack().fillna(0)
                            fig, ax = plt.subplots(figsize=(8, 5))
                            stacked_data.plot(kind='bar', stacked=True, ax=ax, color=colors)
                            ax.set_xlabel(stack_by, fontsize=12)
                            ax.set_ylabel('Count of Responses', fontsize=12)
                            ax.tick_params(axis='x', rotation=0)
                            plt.tight_layout()
                            st.pyplot(fig)
                        else:
                            st.write("No categorical columns available to stack by.")
                        insights[col] = ""
                        raw_texts[col] = ""
                    else:
                        insights[col] = ""
                        raw_texts[col] = ""

                elif data_type == 'Text':
                    if chart_type == 'Word Cloud':
                        text = ' '.join(df[col].dropna().astype(str))
                        if text.strip() == '':
                            st.write("No text data available to generate a word cloud.")
                            fig = None
                        else:
                            vectorizer = CountVectorizer(ngram_range=(1, 3), stop_words='english')
                            text_vector = vectorizer.fit_transform([text])
                            word_freq = dict(zip(vectorizer.get_feature_names_out(), text_vector.toarray().flatten()))
                            sorted_word_freq = dict(sorted(word_freq.items(), key=lambda item: item[1], reverse=True)[:50])

                            wordcloud = WordCloud(
                                width=800, height=400,
                                colormap=selected_colormap, background_color='white',
                                prefer_horizontal=0.8, max_font_size=80, min_font_size=12
                            ).generate_from_frequencies(sorted_word_freq)
                            
                            fig, ax = plt.subplots(figsize=(10, 5))
                            ax.imshow(wordcloud, interpolation='bilinear')
                            ax.axis('off')
                            st.pyplot(fig)

                            total_responses = len(df[col].dropna())
                            if total_responses > 0 and sorted_word_freq:
                                # Get the top 2 phrases from each n-gram level (1, 2, 3)
                                ngram_phrases = {1: [], 2: [], 3: []}
                                for phrase in sorted_word_freq.keys():
                                    ngram_size = len(phrase.split())
                                    if ngram_size in ngram_phrases and len(ngram_phrases[ngram_size]) < 2:
                                        ngram_phrases[ngram_size].append(phrase)

                                top_phrases = [phrase for phrases in ngram_phrases.values() for phrase in phrases]
                                insights[col] = f"Top phrases: {', '.join(top_phrases)}."
                                st.markdown(f"**Insight**: {insights[col]}")
                            else:
                                insights[col] = ""
                                st.markdown("**Insight**: No responses were available for this question.")
                        raw_texts[col] = ""
                    elif chart_type == 'Raw Responses':
                        raw_responses = df[col].dropna().astype(str).tolist()
                        raw_text = "\n\n".join(raw_responses)

                        # Display as table in app
                        st.dataframe(pd.DataFrame({col: raw_responses}))

                        fig = None
                        insights[col] = ""
                        raw_texts[col] = raw_text
                    else:
                        insights[col] = ""
                        raw_texts[col] = ""
                else:
                    insights[col] = ""
                    raw_texts[col] = ""
                    fig = None

                # Store the figure or raw text for exporting
                if fig:
                    charts[col] = fig
                    raw_texts[col] = ""
                elif raw_text:
                    charts[col] = None
                    raw_texts[col] = raw_text
                else:
                    charts[col] = None
                    raw_texts[col] = ""

        # Export to PDF or PPTX
        st.subheader("Export Results")
        export_format = st.selectbox("Select export format", ["PDF", "PPTX"], key='export_format')

        if st.button("Export Results"):
            if export_format == "PDF":
                pdf_bytes = export_to_pdf(charts, insights, raw_texts)
                st.download_button("Download PDF", data=pdf_bytes, file_name="survey_results.pdf", mime="application/pdf")
            elif export_format == "PPTX":
                ppt_bytes = export_to_pptx(charts, insights, raw_texts)
                st.download_button("Download PPTX", data=ppt_bytes, file_name="survey_results.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

        # Display number of rows and columns at the bottom
        st.write(f"Number of rows in dataset: {df.shape[0]}")
        st.write(f"Columns in dataset: {list(df.columns)}")

    except Exception as e:
        st.error(f"An error occurred: {e}")